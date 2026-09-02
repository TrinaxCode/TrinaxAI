"""Text extraction for rich document types (PDF, Word, Excel, etc.).

The agent's ``read_file`` tool falls back to this module when a file is not
plain UTF-8 text, so the agent can read the same document types TrinaxAI indexes
(PDF, .docx, .pptx, .xlsx, .odt, .rtf, .epub, .csv…). Each extractor uses the
standalone parsing libraries that ship with TrinaxAI's requirements and is
imported lazily so a missing optional dependency only disables its own type.

Every function returns plain text or raises; the caller converts failures into a
short error string. Extraction is best-effort and returns visible text only — no
layout, styles or embedded media.
"""

from __future__ import annotations

import zipfile
from html.parser import HTMLParser
from pathlib import Path

from defusedxml import ElementTree

# Extensions handled here rather than as raw UTF-8 text. Kept in sync with the
# document types TrinaxAI's indexer supports.
DOCUMENT_EXTENSIONS = {
    ".pdf",
    ".docx",
    ".doc",
    ".pptx",
    ".ppt",
    ".xlsx",
    ".xls",
    ".odt",
    ".odp",
    ".ods",
    ".rtf",
    ".epub",
}
_EPUB_TEXT_LIMIT = 20 * 1024 * 1024


class _HTMLText(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.parts: list[str] = []

    def handle_data(self, data: str) -> None:
        if data.strip():
            self.parts.append(data.strip())


def is_document(path: Path) -> bool:
    return path.suffix.lower() in DOCUMENT_EXTENSIONS


def _extract_pdf(path: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    pages = []
    for index, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        if text.strip():
            pages.append(f"[page {index}]\n{text.strip()}")
    return "\n\n".join(pages)


def _extract_docx(path: Path) -> str:
    import docx2txt

    return docx2txt.process(str(path)) or ""


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    chunks = []
    for index, slide in enumerate(prs.slides, start=1):
        lines = [shape.text for shape in slide.shapes if getattr(shape, "has_text_frame", False) and shape.text.strip()]
        if lines:
            chunks.append(f"[slide {index}]\n" + "\n".join(lines))
    return "\n\n".join(chunks)


def _extract_xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    chunks = []
    for sheet in wb.worksheets:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                rows.append("\t".join(cells))
        if rows:
            chunks.append(f"[sheet {sheet.title}]\n" + "\n".join(rows))
    wb.close()
    return "\n\n".join(chunks)


def _extract_odf(path: Path) -> str:
    from odf import teletype, text
    from odf.opendocument import load

    doc = load(str(path))
    paragraphs = [teletype.extractText(node) for node in doc.getElementsByType(text.P)]
    return "\n".join(p for p in paragraphs if p.strip())


def _extract_rtf(path: Path) -> str:
    from striprtf.striprtf import rtf_to_text

    return rtf_to_text(path.read_text(encoding="utf-8", errors="replace"))


def _extract_epub(path: Path) -> str:
    chunks = []
    extracted_bytes = 0
    with zipfile.ZipFile(path) as archive:
        for info in archive.infolist():
            if Path(info.filename).suffix.lower() not in {".html", ".htm", ".xhtml", ".xml", ".ncx"}:
                continue
            if info.file_size > _EPUB_TEXT_LIMIT:
                continue
            extracted_bytes += info.file_size
            if extracted_bytes > _EPUB_TEXT_LIMIT:
                raise ValueError("EPUB expanded text exceeds the safe extraction limit")
            source = archive.read(info)
            try:
                root = ElementTree.fromstring(source)
                text = "\n".join(value.strip() for value in root.itertext() if value.strip())
            except ElementTree.ParseError:
                parser = _HTMLText()
                parser.feed(source.decode("utf-8", errors="replace"))
                text = "\n".join(parser.parts)
            if text:
                chunks.append(f"[{info.filename}]\n{text}")
    return "\n\n".join(chunks)


_EXTRACTORS = {
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".doc": _extract_docx,
    ".pptx": _extract_pptx,
    ".ppt": _extract_pptx,
    ".xlsx": _extract_xlsx,
    ".xls": _extract_xlsx,
    ".odt": _extract_odf,
    ".odp": _extract_odf,
    ".ods": _extract_odf,
    ".rtf": _extract_rtf,
    ".epub": _extract_epub,
}


def extract_document_text(path: Path) -> str:
    """Extract visible text from a supported document, or raise on failure."""
    extractor = _EXTRACTORS.get(path.suffix.lower())
    if extractor is None:
        raise ValueError(f"unsupported document type: {path.suffix}")
    return (extractor(path) or "").strip()
