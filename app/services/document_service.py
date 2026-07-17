"""Document extraction services."""

from __future__ import annotations

# ruff: noqa: F405
from .shared_runtime import *  # noqa: F403


def _decode_text_bytes(data: bytes) -> str:
    for encoding in ("utf-8", "utf-8-sig", "latin-1"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def _extract_pdf_text(data: bytes) -> str:
    try:
        from pypdf import PdfReader
    except Exception as exc:
        raise HTTPException(status_code=501, detail="PDF extraction requires pypdf.") from exc
    try:
        reader = PdfReader(BytesIO(data))
        pages: list[str] = []
        for index, page in enumerate(reader.pages, start=1):
            text = page.extract_text() or ""
            if text.strip():
                pages.append(f"[Page {index}]\n{text.strip()}")
        text_result = "\n\n".join(pages).strip()
        # OCR fallback (Phase 5.1): when pypdf returns very little text (scanned PDF)
        # and the user has enabled OCR via TRINAXAI_OCR=1, rasterize the pages and
        # run tesseract. Failures degrade gracefully — we just keep the original text.
        if config.TRINAXAI_OCR and len(text_result) < 50:
            try:
                import pytesseract  # type: ignore
                from pdf2image import convert_from_bytes  # type: ignore

                images = convert_from_bytes(data, dpi=200)
                ocr_pages: list[str] = []
                for i, img in enumerate(images, start=1):
                    ocr_text = pytesseract.image_to_string(img, lang="eng+spa") or ""
                    if ocr_text.strip():
                        ocr_pages.append(f"[Page {i}]\n{ocr_text.strip()}")
                ocr_result = "\n\n".join(ocr_pages).strip()
                if ocr_result and len(ocr_result) > len(text_result):
                    return ocr_result
            except Exception:
                # OCR not installed or failed; fall through to original text.
                pass
        return text_result
    except Exception as exc:
        raise HTTPException(status_code=422, detail=f"Could not extract PDF text: {str(exc)[:180]}") from exc


def _extract_docx_text(data: bytes) -> str:
    try:
        import docx2txt
    except Exception as exc:
        raise HTTPException(status_code=501, detail="DOCX extraction requires docx2txt.") from exc
    tmp_path = ""
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as tmp:
            tmp.write(data)
            tmp_path = tmp.name
        return (docx2txt.process(tmp_path) or "").strip()
    except Exception as exc:
        raise HTTPException(status_code=422, detail=f"Could not extract DOCX text: {str(exc)[:180]}") from exc
    finally:
        if tmp_path:
            try:
                os.remove(tmp_path)
            except OSError:
                pass


def _extract_pptx_text(data: bytes) -> str:
    try:
        from pptx import Presentation
    except Exception as exc:
        raise HTTPException(status_code=501, detail="PPTX extraction requires python-pptx.") from exc
    try:
        presentation = Presentation(BytesIO(data))
        slides: list[str] = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            parts: list[str] = []
            for shape in slide.shapes:
                text = ""
                if getattr(shape, "has_text_frame", False):
                    text = shape.text or ""
                elif getattr(shape, "has_table", False):
                    rows: list[str] = []
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if cells:
                            rows.append(" | ".join(cells))
                    text = "\n".join(rows)
                if text.strip():
                    parts.append(text.strip())
            try:
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text or ""
                    if notes.strip():
                        parts.append(f"Notes:\n{notes.strip()}")
            except Exception:
                LOG.debug("Best-effort operation failed", exc_info=True)
            if parts:
                slides.append(f"[Slide {slide_index}]\n" + "\n\n".join(parts))
        return "\n\n".join(slides).strip()
    except Exception as exc:
        raise HTTPException(status_code=422, detail=f"Could not extract PPTX text: {str(exc)[:180]}") from exc


def _extract_xlsx_text(data: bytes) -> str:
    try:
        from openpyxl import load_workbook
    except Exception as exc:
        raise HTTPException(status_code=501, detail="Spreadsheet extraction requires openpyxl.") from exc
    try:
        workbook = load_workbook(BytesIO(data), read_only=True, data_only=True)
        sheets: list[str] = []
        for sheet in workbook.worksheets:
            rows: list[str] = []
            for row in sheet.iter_rows(values_only=True):
                values = [str(value).strip() if value is not None else "" for value in row]
                while values and not values[-1]:
                    values.pop()
                if any(values):
                    rows.append("\t".join(values))
            if rows:
                sheets.append(f"[Sheet: {sheet.title}]\n" + "\n".join(rows))
        workbook.close()
        return "\n\n".join(sheets).strip()
    except Exception as exc:
        raise HTTPException(status_code=422, detail=f"Could not extract spreadsheet text: {str(exc)[:180]}") from exc


def _extract_rtf_text(data: bytes) -> str:
    try:
        from striprtf.striprtf import rtf_to_text
    except Exception as exc:
        raise HTTPException(status_code=501, detail="RTF extraction requires striprtf.") from exc
    try:
        return rtf_to_text(_decode_text_bytes(data)).strip()
    except Exception as exc:
        raise HTTPException(status_code=422, detail=f"Could not extract RTF text: {str(exc)[:180]}") from exc


def _extract_odf_text(data: bytes) -> str:
    try:
        from odf import teletype
        from odf.opendocument import load
        from odf.table import Table, TableCell, TableRow
        from odf.text import H, P
    except Exception as exc:
        raise HTTPException(status_code=501, detail="OpenDocument extraction requires odfpy.") from exc
    try:
        document = load(BytesIO(data))
        parts: list[str] = []
        for node_type in (H, P):
            for node in document.getElementsByType(node_type):
                text = teletype.extractText(node).strip()
                if text:
                    parts.append(text)
        for table in document.getElementsByType(Table):
            rows: list[str] = []
            for row in table.getElementsByType(TableRow):
                cells = [teletype.extractText(cell).strip() for cell in row.getElementsByType(TableCell)]
                while cells and not cells[-1]:
                    cells.pop()
                if any(cells):
                    rows.append("\t".join(cells))
            if rows:
                parts.append("\n".join(rows))
        return "\n".join(parts).strip()
    except Exception as exc:
        raise HTTPException(status_code=422, detail=f"Could not extract OpenDocument text: {str(exc)[:180]}") from exc


def _convert_office_bytes(data: bytes, source_ext: str, target_ext: str) -> bytes:
    """Convert legacy/OpenDocument containers through local LibreOffice."""
    executable = shutil.which("libreoffice") or shutil.which("soffice")
    if not executable:
        raise HTTPException(
            status_code=501,
            detail=f"{source_ext.upper()} extraction requires LibreOffice.",
        )
    with tempfile.TemporaryDirectory(prefix="trinaxai-office-") as directory:
        source = os.path.join(directory, f"document{source_ext}")
        with open(source, "wb") as stream:
            stream.write(data)
        try:
            result = subprocess.run(
                [executable, "--headless", "--convert-to", target_ext.lstrip("."), "--outdir", directory, source],
                capture_output=True,
                text=True,
                timeout=90,
            )
        except (OSError, subprocess.TimeoutExpired) as exc:
            raise HTTPException(status_code=422, detail=f"Office conversion failed: {str(exc)[:180]}") from exc
        target = os.path.join(directory, f"document{target_ext}")
        if result.returncode != 0 or not os.path.isfile(target):
            detail = (result.stderr or result.stdout or "conversion produced no output").strip()
            raise HTTPException(status_code=422, detail=f"Office conversion failed: {detail[:180]}")
        with open(target, "rb") as stream:
            return stream.read()


def _extract_document_text(filename: str, data: bytes) -> str:
    ext = os.path.splitext(filename.lower())[1]
    if ext == ".pdf":
        return _extract_pdf_text(data)
    if ext == ".docx":
        return _extract_docx_text(data)
    if ext == ".pptx":
        return _extract_pptx_text(data)
    if ext == ".xlsx":
        return _extract_xlsx_text(data)
    if ext == ".rtf":
        return _extract_rtf_text(data)
    if ext in {".odt", ".ods", ".odp"}:
        return _extract_odf_text(data)
    if ext == ".doc":
        return _extract_docx_text(_convert_office_bytes(data, ext, ".docx"))
    if ext == ".ppt":
        return _extract_pptx_text(_convert_office_bytes(data, ext, ".pptx"))
    if ext == ".xls":
        return _extract_xlsx_text(_convert_office_bytes(data, ext, ".xlsx"))
    if ext in {
        ".txt",
        ".md",
        ".mdx",
        ".rst",
        ".csv",
        ".json",
        ".xml",
        ".yml",
        ".yaml",
        ".toml",
        ".ini",
        ".log",
    }:
        return _decode_text_bytes(data).strip()
    return _decode_text_bytes(data).strip()


async def document_extract(request: Request, file: UploadFile = File(...)):
    """Extract plain text from an uploaded document (PDF/DOCX/PPTX/text).

    Extrae texto plano de un documento subido (PDF/DOCX/PPTX/texto) sin
    indexarlo. Devuelve el texto (posiblemente truncado) para uso puntual.
    """
    enforce_rate_limit(request, bucket="document_extract")
    name = file.filename or "document"
    data = bytearray()
    total_bytes = 0
    while True:
        chunk = await file.read(1024 * 1024)
        if not chunk:
            break
        total_bytes += len(chunk)
        if total_bytes > DOC_EXTRACT_MAX_BYTES:
            await file.close()
            raise HTTPException(
                status_code=413,
                detail=(f"Document is too large for temporary extraction. Limit: {DOC_EXTRACT_MAX_BYTES} bytes."),
            )
        data.extend(chunk)
    await file.close()
    if not data:
        raise HTTPException(status_code=400, detail="Empty file.")

    def extract_with_slot():
        with _document_slots:
            return _extract_document_text(name, data)

    text = await run_in_threadpool(extract_with_slot)
    if not text.strip():
        raise HTTPException(status_code=422, detail="No readable text found in this document.")
    truncated = len(text) > DOC_EXTRACT_MAX_CHARS
    if truncated:
        text = text[:DOC_EXTRACT_MAX_CHARS]
    return {
        "ok": True,
        "name": name,
        "text": text,
        "chars": len(text),
        "truncated": truncated,
    }


__all__ = [name for name in globals() if not name.startswith("__")]
