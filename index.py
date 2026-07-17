"""
TrinaxAI — Indexador de documentos.

Características:
  • Chunking consciente del lenguaje: CodeSplitter (AST) para código,
    SentenceSplitter para prosa. No parte funciones por la mitad.
  • Embeddings bge-m3 (multilingüe, 1024 dims).
  • Metadata de proyecto en cada chunk (para citas y filtro por proyecto).
  • INCREMENTAL: solo re-indexa archivos nuevos o modificados (fingerprint
    de contenido + versión de pipeline). Actualizar = segundos, no horas.
  • Publicación recuperable: índice y manifiesto cambian como una generación,
    con rollback automático si el proceso se interrumpe.
  • Múltiples raíces por colección sin colisiones ni borrados cruzados.
  • Sin LLM cargado al indexar (solo hace falta el embedder).
"""

from __future__ import annotations

import hashlib
import json
import os
import shutil
import subprocess
import sys
import tempfile
import zipfile
from concurrent.futures import ThreadPoolExecutor
from dataclasses import dataclass, field
from email import policy
from email.parser import BytesParser
from html.parser import HTMLParser
from typing import TYPE_CHECKING

if TYPE_CHECKING:
    from llama_index.core import VectorStoreIndex
    from llama_index.core.schema import Document

# On Windows, stdout defaults to cp1252 which can't encode emoji/Unicode.
# Wrap it so the indexer doesn't crash mid-job on a harmless print.
if sys.platform == "win32" and hasattr(sys.stdout, "buffer"):
    import codecs

    sys.stdout = codecs.getwriter("utf-8")(sys.stdout.buffer, "replace")  # type: ignore[assignment]

import config
from trinaxai_core import (
    _positive_int,
    exclusive_process_lock,
    sanitize_collection_id,
    source_id_for_root,
)
from trinaxai_index_storage import (
    atomic_write_json,
    publish_index_generation,
    recover_interrupted_transaction,
)

EXTRACTOR_EXTS = {".pdf", ".docx"}
TEXT_ENCODINGS = ("utf-8-sig", "utf-8", "cp1252", "latin-1")
_TEXT_SAMPLE_BYTES = 8192
_EPUB_TEXT_LIMIT = 20 * 1024 * 1024
MANIFEST_SCHEMA_VERSION = 2
FINGERPRINT_ALGORITHM = "blake2b-256"
_HASH_BLOCK_BYTES = 1024 * 1024


class _HTMLTextExtractor(HTMLParser):
    """Small dependency-free HTML-to-text extractor."""

    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.parts: list[str] = []
        self._ignored_depth = 0

    def handle_starttag(self, tag: str, attrs) -> None:
        if tag in {"script", "style", "noscript", "svg"}:
            self._ignored_depth += 1
        elif tag in {"p", "div", "br", "li", "h1", "h2", "h3", "h4", "tr"}:
            self.parts.append("\n")

    def handle_endtag(self, tag: str) -> None:
        if tag in {"script", "style", "noscript", "svg"} and self._ignored_depth:
            self._ignored_depth -= 1
        elif tag in {"p", "div", "li", "h1", "h2", "h3", "h4", "tr"}:
            self.parts.append("\n")

    def handle_data(self, data: str) -> None:
        if not self._ignored_depth and data.strip():
            self.parts.append(data.strip())

    def text(self) -> str:
        lines = (" ".join(line.split()) for line in " ".join(self.parts).splitlines())
        return "\n".join(line for line in lines if line).strip()


def _html_to_text(value: str) -> str:
    parser = _HTMLTextExtractor()
    parser.feed(value)
    parser.close()
    return parser.text()


def _env_int(name: str, default: int, *, minimum: int = 1, maximum: int | None = None) -> int:
    """Read an int env var, clamped to [minimum, maximum]; falls back on bad input.

    Thin wrapper over ``trinaxai_core._positive_int`` so parsing/clamping stays
    consistent across config/index/core.
    """
    return _positive_int(os.getenv(name, default), default, minimum=minimum, maximum=maximum)


# ==================== SETTINGS ====================
# NO se define Settings.llm: indexar solo necesita embeddings.
COLLECTION_ID = sanitize_collection_id(
    os.getenv("TRINAXAI_COLLECTION_ID", config.DEFAULT_COLLECTION_ID),
    fallback=config.DEFAULT_COLLECTION_ID,
)
COLLECTION_NAME = (
    os.getenv("TRINAXAI_COLLECTION_NAME", config.DEFAULT_COLLECTION_NAME).strip() or config.DEFAULT_COLLECTION_NAME
)

INDEX_BATCH_SIZE = _env_int("TRINAXAI_INDEX_BATCH_SIZE", 100, minimum=1, maximum=1000)
INDEX_NODE_BATCH_SIZE = _env_int("TRINAXAI_INDEX_NODE_BATCH_SIZE", 32, minimum=1, maximum=256)
INDEX_LOAD_WORKERS = _env_int(
    "TRINAXAI_INDEX_LOAD_WORKERS",
    min(8, os.cpu_count() or 4),
    minimum=1,
    maximum=32,
)
APPEND_ONLY = os.getenv("TRINAXAI_INDEX_APPEND", "0").strip().lower() in {
    "1",
    "true",
    "yes",
    "on",
}


@dataclass(frozen=True)
class SourceContext:
    """Identity and path policy for one independently synchronized source root."""

    root: str
    source_id: str
    collection_id: str
    collection_name: str
    project_name: str

    @classmethod
    def create(
        cls,
        root: str,
        *,
        source_id: str | None = None,
        collection_id: str | None = None,
        collection_name: str | None = None,
    ) -> "SourceContext":
        canonical_root = os.path.realpath(os.path.abspath(os.path.expanduser(root)))
        basename = os.path.basename(canonical_root.rstrip(os.sep)) or "root"
        explicit_id = source_id or os.getenv("TRINAXAI_SOURCE_ID")
        safe_source_id = source_id_for_root(canonical_root, explicit_id=explicit_id)
        return cls(
            root=canonical_root,
            source_id=safe_source_id,
            collection_id=sanitize_collection_id(
                collection_id or COLLECTION_ID,
                fallback=config.DEFAULT_COLLECTION_ID,
            ),
            collection_name=(collection_name or COLLECTION_NAME).strip() or config.DEFAULT_COLLECTION_NAME,
            project_name=basename,
        )

    def relative_path(self, path: str) -> str:
        absolute = os.path.realpath(os.path.abspath(path))
        try:
            if os.path.commonpath([self.root, absolute]) != self.root:
                raise ValueError(f"Path is outside source root: {path}")
        except ValueError as exc:
            raise ValueError(f"Path is outside source root: {path}") from exc
        relative = os.path.relpath(absolute, self.root).replace("\\", "/")
        return relative

    def source_key_for_relative(self, relative: str) -> str:
        clean_relative = relative.replace("\\", "/").lstrip("/")
        return f"{self.collection_id}:{self.source_id}:{clean_relative}"

    def source_key(self, path: str) -> str:
        return self.source_key_for_relative(self.relative_path(path))


def _default_source_context() -> SourceContext:
    return SourceContext.create(config.PROJECTS_DIRS[0])


def _pipeline_version() -> str:
    """Stable version of every setting that changes generated chunks/vectors."""
    inputs = {
        "schema": MANIFEST_SCHEMA_VERSION,
        "embed_model": config.EMBED_MODEL,
        "embed_dims": config.EMBED_DIMS,
        "chunk_size": config.CHUNK_SIZE,
        "chunk_overlap": config.CHUNK_OVERLAP,
        "code_chunk_lines": config.CODE_CHUNK_LINES,
        "code_chunk_overlap": config.CODE_CHUNK_LINES_OVERLAP,
        "code_max_chars": config.CODE_MAX_CHARS,
        "extractor_version": 2,
    }
    payload = json.dumps(inputs, sort_keys=True, separators=(",", ":")).encode("utf-8")
    return hashlib.sha256(payload).hexdigest()[:20]

# Cache de CodeSplitters por lenguaje (crearlos es caro).
_code_splitters: dict[str, object] = {}
_prose_splitter = None
_embed_configured = False


def ensure_embed_settings() -> None:
    """Initialize Ollama embeddings only when an actual index run starts."""
    global _embed_configured
    if not _embed_configured:
        from llama_index.core import Settings

        Settings.embed_model = config.make_embed()
        _embed_configured = True


def _code_splitter(language: str):
    if language not in _code_splitters:
        from llama_index.core.node_parser import CodeSplitter

        _code_splitters[language] = CodeSplitter(
            language=language,
            chunk_lines=config.CODE_CHUNK_LINES,
            chunk_lines_overlap=config.CODE_CHUNK_LINES_OVERLAP,
            max_chars=config.CODE_MAX_CHARS,
        )
    return _code_splitters[language]


def _sentence_splitter():
    global _prose_splitter
    if _prose_splitter is None:
        from llama_index.core.node_parser import SentenceSplitter

        _prose_splitter = SentenceSplitter(
            chunk_size=config.CHUNK_SIZE,
            chunk_overlap=config.CHUNK_OVERLAP,
        )
    return _prose_splitter


# ==================== LECTOR DE ARCHIVOS ====================
def collect_files(root: str) -> list[str]:
    """Recorre `root` PODANDO carpetas de dependencias.

    No desciende en node_modules/site-packages/venv/etc., así que es
    órdenes de magnitud más rápido que enumerar todo y filtrar después.
    """
    allowed = {e.lower() for e in config.REQUIRED_EXTS}
    allowed_names = {
        "dockerfile",
        "makefile",
        "readme",
        "license",
        "changelog",
        "contributing",
        "gemfile",
        "procfile",
        ".env",
    }
    files: list[str] = []
    skipped_big = 0
    file_count = 0
    for dirpath, dirnames, filenames in os.walk(root):
        # Podar: quita in-place las carpetas excluidas y las ocultas.
        dirnames[:] = [
            d
            for d in dirnames
            if d not in config.EXCLUDE_DIR_NAMES
            and not d.startswith(".")
            and not d.endswith((".egg-info", ".dist-info"))
        ]
        for fn in filenames:
            if fn.startswith(".") and fn.lower() not in allowed_names:
                continue
            full = os.path.join(dirpath, fn)
            if os.path.islink(full):
                continue
            try:
                if os.path.getsize(full) > config.max_file_bytes(full):
                    skipped_big += 1
                    continue
            except OSError:
                continue
            known_type = fn.lower() in allowed_names or os.path.splitext(fn)[1].lower() in allowed
            # Unknown extensions are still indexed when their bytes look like
            # human-readable text. This covers domain-specific formats without
            # blindly feeding executables, archives or media to the embedder.
            if not known_type and not _is_probably_text_file(full):
                continue
            files.append(full)
            file_count += 1
            if file_count % 5000 == 0:
                print(f"   📂 {file_count} archivos encontrados...", flush=True)
    if file_count >= 5000:
        print(f"   📂 {file_count} archivos encontrados en total")
    if skipped_big:
        print(f"   ⏭️  {skipped_big} archivos omitidos por tamaño (sobre el límite configurado para su tipo de archivo)")
    return files


def _is_probably_text_file(path: str) -> bool:
    """Detect textual files independently of their extension."""
    try:
        with open(path, "rb") as stream:
            sample = stream.read(_TEXT_SAMPLE_BYTES)
    except OSError:
        return False
    if not sample:
        return True
    if b"\x00" in sample:
        return False
    try:
        decoded = sample.decode("utf-8")
    except UnicodeDecodeError:
        decoded = sample.decode("cp1252", errors="replace")
    controls = sum(1 for char in decoded if ord(char) < 32 and char not in {"\n", "\r", "\t", "\f", "\b"})
    replacements = decoded.count("\ufffd")
    return (controls + replacements) / max(1, len(decoded)) < 0.02


def _rel(path: str, context: SourceContext | None = None) -> str:
    """Return a source-relative path (legacy-compatible public helper)."""
    if context is not None:
        return context.relative_path(path)
    try:
        return os.path.relpath(path, config.PROJECTS_DIRS[0]).replace("\\", "/")
    except ValueError:
        return path.replace("\\", "/")


def _source_key(path: str, context: SourceContext | None = None) -> str:
    """Return a source-aware key; retain the legacy shape without context."""
    if context is not None:
        return context.source_key(path)
    return f"{COLLECTION_ID}:{_rel(path)}"


def _decode_text_bytes(data: bytes) -> str:
    """Decode source files without depending on the OS locale.

    Windows often defaults to cp1252, which can raise ``charmap`` errors on
    bytes that are valid in other encodings. Latin-1 is the final fallback
    because it maps every byte and keeps indexing from aborting.
    """
    if data.startswith((b"\xff\xfe", b"\xfe\xff")) or b"\x00" in data[:200]:
        try:
            return data.decode("utf-16")
        except UnicodeDecodeError:
            pass
    for encoding in TEXT_ENCODINGS:
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def _load_text_document(path: str) -> Document:
    from llama_index.core.schema import Document

    with open(path, "rb") as f:
        text = _decode_text_bytes(f.read())
    return Document(text=text, metadata={"file_path": path})


def _document(path: str, text: str) -> Document:
    from llama_index.core.schema import Document

    cleaned = text.strip()
    if not cleaned:
        raise ValueError("the file contains no extractable text")
    return Document(text=cleaned, metadata={"file_path": path})


def _load_html_document(path: str) -> Document:
    with open(path, "rb") as stream:
        return _document(path, _html_to_text(_decode_text_bytes(stream.read())))


def _load_notebook_document(path: str) -> Document:
    with open(path, encoding="utf-8") as stream:
        notebook = json.load(stream)
    parts: list[str] = []
    for index_number, cell in enumerate(notebook.get("cells") or [], start=1):
        if not isinstance(cell, dict):
            continue
        cell_type = str(cell.get("cell_type") or "cell")
        source = cell.get("source") or ""
        text = "".join(source) if isinstance(source, list) else str(source)
        if text.strip():
            parts.append(f"[{cell_type.title()} cell {index_number}]\n{text.strip()}")
    return _document(path, "\n\n".join(parts))


def _load_email_document(path: str) -> Document:
    with open(path, "rb") as stream:
        message = BytesParser(policy=policy.default).parse(stream)
    headers = [f"{name}: {message.get(name)}" for name in ("Subject", "From", "To", "Date") if message.get(name)]
    bodies: list[str] = []
    parts = message.walk() if message.is_multipart() else [message]
    for part in parts:
        if part.get_content_disposition() == "attachment":
            continue
        content_type = part.get_content_type()
        if content_type not in {"text/plain", "text/html"}:
            continue
        try:
            value = part.get_content()
        except Exception:
            payload = part.get_payload(decode=True) or b""
            value = _decode_text_bytes(payload)
        bodies.append(_html_to_text(value) if content_type == "text/html" else value.strip())
    return _document(path, "\n".join(headers) + "\n\n" + "\n\n".join(bodies))


def _load_epub_document(path: str) -> Document:
    sections: list[str] = []
    extracted_bytes = 0
    with zipfile.ZipFile(path) as archive:
        for info in archive.infolist():
            ext = os.path.splitext(info.filename.lower())[1]
            if ext not in {".html", ".htm", ".xhtml", ".xml", ".ncx"}:
                continue
            if info.file_size > _EPUB_TEXT_LIMIT:
                continue
            extracted_bytes += info.file_size
            if extracted_bytes > _EPUB_TEXT_LIMIT:
                raise ValueError("EPUB expanded text exceeds the safe extraction limit")
            value = _html_to_text(_decode_text_bytes(archive.read(info)))
            if value:
                sections.append(f"[{info.filename}]\n{value}")
    return _document(path, "\n\n".join(sections))


def _load_extracted_documents(path: str) -> list[Document]:
    from llama_index.core import SimpleDirectoryReader

    return SimpleDirectoryReader(
        input_files=[path],
        exclude_hidden=False,
        exclude=config.EXCLUDE_PATTERNS,
        encoding="utf-8",
        errors="replace",
        raise_on_error=False,
    ).load_data()


def emit_progress(phase: str, **values: object) -> None:
    print("TRINAXAI_PROGRESS " + json.dumps({"phase": phase, **values}, ensure_ascii=False), flush=True)


def _load_pdf_documents(path: str) -> list[Document]:
    """Extract a PDF page-by-page with observable, bounded work units."""
    from llama_index.core.schema import Document
    from pypdf import PdfReader

    reader = PdfReader(path, strict=False)
    total = len(reader.pages)
    emit_progress("extracting", pages_total=total, pages_processed=0, determinate=bool(total))
    documents: list[Document] = []
    for page_number, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        if text:
            documents.append(Document(text=f"[Page {page_number}]\n{text}", metadata={"file_path": path, "page": page_number}))
        emit_progress("extracting", pages_total=total, pages_processed=page_number, determinate=True)
    if not documents:
        raise ValueError("PDF contains no extractable text; OCR may be required")
    return documents


def _load_pptx_document(path: str) -> Document:
    from llama_index.core.schema import Document
    from pptx import Presentation

    presentation = Presentation(path)
    slides: list[str] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            text = ""
            if getattr(shape, "has_text_frame", False):
                text = shape.text or ""
            elif getattr(shape, "has_table", False):
                rows: list[str] = []
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        rows.append(" | ".join(cells))
                text = "\n".join(rows)
            if text.strip():
                parts.append(text.strip())
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text or ""
                if notes.strip():
                    parts.append(f"Notes:\n{notes.strip()}")
        except Exception:
            pass
        if parts:
            slides.append(f"[Slide {slide_index}]\n" + "\n\n".join(parts))
    return Document(text="\n\n".join(slides).strip(), metadata={"file_path": path})


def _load_xlsx_document(path: str) -> Document:
    from llama_index.core.schema import Document
    from openpyxl import load_workbook

    workbook = load_workbook(path, read_only=True, data_only=True)
    sheets: list[str] = []
    try:
        for sheet in workbook.worksheets:
            rows: list[str] = []
            for row in sheet.iter_rows(values_only=True):
                values = [str(value).strip() if value is not None else "" for value in row]
                while values and not values[-1]:
                    values.pop()
                if any(values):
                    rows.append("\t".join(values))
            if rows:
                sheets.append(f"[Sheet: {sheet.title}]\n" + "\n".join(rows))
    finally:
        workbook.close()
    return Document(text="\n\n".join(sheets).strip(), metadata={"file_path": path})


def _load_rtf_document(path: str) -> Document:
    from llama_index.core.schema import Document
    from striprtf.striprtf import rtf_to_text

    with open(path, "rb") as stream:
        text = rtf_to_text(_decode_text_bytes(stream.read()))
    return Document(text=text.strip(), metadata={"file_path": path})


def _load_odf_document(path: str) -> Document:
    from llama_index.core.schema import Document
    from odf import teletype
    from odf.opendocument import load
    from odf.table import Table, TableCell, TableRow
    from odf.text import H, P

    document = load(path)
    parts: list[str] = []
    for node_type in (H, P):
        for node in document.getElementsByType(node_type):
            text = teletype.extractText(node).strip()
            if text:
                parts.append(text)
    for table in document.getElementsByType(Table):
        rows: list[str] = []
        for row in table.getElementsByType(TableRow):
            cells = [teletype.extractText(cell).strip() for cell in row.getElementsByType(TableCell)]
            while cells and not cells[-1]:
                cells.pop()
            if any(cells):
                rows.append("\t".join(cells))
        if rows:
            parts.append("\n".join(rows))
    return Document(text="\n".join(parts).strip(), metadata={"file_path": path})


def _load_converted_office_document(path: str, target_ext: str) -> list[Document]:
    executable = shutil.which("libreoffice") or shutil.which("soffice")
    if not executable:
        raise RuntimeError(f"{os.path.splitext(path)[1].upper()} requires LibreOffice")
    with tempfile.TemporaryDirectory(prefix="trinaxai-index-office-") as directory:
        result = subprocess.run(
            [executable, "--headless", "--convert-to", target_ext.lstrip("."), "--outdir", directory, path],
            capture_output=True,
            text=True,
            timeout=90,
        )
        target = os.path.join(directory, f"{os.path.splitext(os.path.basename(path))[0]}{target_ext}")
        if result.returncode != 0 or not os.path.isfile(target):
            detail = (result.stderr or result.stdout or "conversion produced no output").strip()
            raise RuntimeError(f"Office conversion failed: {detail[:180]}")
        documents = _load_file_documents(target)
        for document in documents:
            document.metadata["file_path"] = path
        return documents


def _load_file_documents(path: str) -> list[Document]:
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf":
        return _load_pdf_documents(path)
    if ext in {".html", ".htm", ".xhtml"}:
        return [_load_html_document(path)]
    if ext == ".ipynb":
        return [_load_notebook_document(path)]
    if ext == ".eml":
        return [_load_email_document(path)]
    if ext == ".epub":
        return [_load_epub_document(path)]
    if ext == ".pptx":
        return [_load_pptx_document(path)]
    if ext == ".xlsx":
        return [_load_xlsx_document(path)]
    if ext == ".rtf":
        return [_load_rtf_document(path)]
    if ext in {".odt", ".ods", ".odp"}:
        return [_load_odf_document(path)]
    if ext == ".doc":
        return _load_converted_office_document(path, ".docx")
    if ext == ".ppt":
        return _load_converted_office_document(path, ".pptx")
    if ext == ".xls":
        return _load_converted_office_document(path, ".xlsx")
    if ext in EXTRACTOR_EXTS:
        return _load_extracted_documents(path)
    return [_load_text_document(path)]


def _load_file_documents_result(path: str) -> tuple[str, list[Document], Exception | None]:
    try:
        return path, _load_file_documents(path), None
    except Exception as exc:
        return path, [], exc


def iter_batches(items: list[str], batch_size: int = INDEX_BATCH_SIZE):
    """Yield stable batches without copying the full indexing workload again."""
    for batch_start in range(0, len(items), batch_size):
        yield items[batch_start : batch_start + batch_size]


def total_batches(items: list[str], batch_size: int = INDEX_BATCH_SIZE) -> int:
    """Number of batches :func:`iter_batches` will yield for ``items``."""
    return (len(items) + batch_size - 1) // batch_size if items else 0


def _emit_embed_progress(done: int, total: int) -> None:
    """Emit a machine-parseable, newline-terminated embedding-progress line.

    tqdm's ``show_progress`` bar uses carriage returns, so the supervising
    ``system_service`` process never sees a new stdout line and the UI bar stalls
    at the first "embedding" hit. Printing one real line per batch (with an
    explicit ``N/M``) lets the supervisor map progress proportionally.
    """
    if total <= 0:
        return
    print(f"🔨 Embeddings lote {done}/{total}...", flush=True)
    emit_progress("embedding", batches_processed=done, batches_total=total, determinate=True)


def insert_node_batches(index, nodes: list, *, initialize: bool = False):
    """Insert bounded batches; progress advances only after a completed batch."""
    from llama_index.core import VectorStoreIndex

    batches = list(iter_batches(nodes, INDEX_NODE_BATCH_SIZE))
    current = index
    for batch_number, batch in enumerate(batches, start=1):
        if current is None and initialize:
            current = VectorStoreIndex(batch, show_progress=False)
        else:
            current.insert_nodes(batch, show_progress=False)
        _emit_embed_progress(batch_number, len(batches))
    return current


@dataclass
class LoadResult:
    documents: list[Document] = field(default_factory=list)
    loaded_paths: set[str] = field(default_factory=set)
    failures: dict[str, str] = field(default_factory=dict)


@dataclass
class PreparedBatch:
    nodes: list = field(default_factory=list)
    indexed_paths: set[str] = field(default_factory=set)
    failures: dict[str, str] = field(default_factory=dict)


def load_docs_with_status(paths: list[str], context: SourceContext | None = None) -> LoadResult:
    """Carga documentos y les pone metadata limpia (proyecto, ruta, archivo).

    - doc.id_ = ruta relativa (ID estable → permite borrado/reinserción
      incremental por archivo).
    - Procesa en batches de 100 para no saturar la memoria con directorios
      muy grandes.
    """
    result = LoadResult()
    source_context = context or _default_source_context()
    if not paths:
        return result
    for batch in iter_batches(paths):
        loaded_results: list[tuple[str, list[Document], Exception | None]] = []
        if INDEX_LOAD_WORKERS > 1 and len(batch) > 1:
            workers = min(INDEX_LOAD_WORKERS, len(batch))
            with ThreadPoolExecutor(max_workers=workers) as executor:
                loaded_results.extend(executor.map(_load_file_documents_result, batch))
        else:
            for path in batch:
                loaded_results.append(_load_file_documents_result(path))
        for fp, group, error in loaded_results:
            if error is not None:
                result.failures[fp] = str(error)[:300]
                print(f"   ⚠️  Error leyendo {os.path.basename(fp)}, se reintentará: {error}")
                continue
            group = [document for document in group if str(document.text or "").strip()]
            if not group:
                result.failures[fp] = "no extractable text"
                print(f"   ⚠️  {os.path.basename(fp)} no contiene texto extraíble; se reintentará")
                continue
            rel = source_context.relative_path(fp)
            document_id = source_context.source_key_for_relative(rel)
            for i, d in enumerate(group):
                d.id_ = document_id if len(group) == 1 else f"{document_id}#{i}"
                metadata = dict(d.metadata or {})
                metadata.update(
                    {
                        "project": source_context.project_name,
                        "rel_path": rel,
                        "file_name": os.path.basename(fp),
                        "source_key": source_context.source_key_for_relative(rel),
                        "source_id": source_context.source_id,
                        "source_root": source_context.root,
                        "collection_id": source_context.collection_id,
                        "collection_name": source_context.collection_name,
                        "pipeline_version": _pipeline_version(),
                    }
                )
                d.metadata = metadata
                d.excluded_embed_metadata_keys = []
                d.excluded_llm_metadata_keys = []
                result.documents.append(d)
            result.loaded_paths.add(fp)
    return result


def load_docs(paths: list[str], context: SourceContext | None = None) -> list[Document]:
    """Compatibility wrapper returning only successfully loaded documents."""
    return load_docs_with_status(paths, context).documents


def build_nodes(documents: list[Document]) -> list:
    """Trocea por extensión: código → AST, prosa → texto. La metadata del
    documento (proyecto, ruta) se hereda automáticamente en cada chunk."""
    nodes = []
    code_count = prose_count = fallback = 0
    for doc in documents:
        file_path = doc.metadata.get("rel_path", "")
        ext = os.path.splitext(file_path)[1].lower()
        language = config.CODE_LANG_BY_EXT.get(ext)

        doc_nodes = None
        if language:
            try:
                doc_nodes = _code_splitter(language).get_nodes_from_documents([doc])
                code_count += 1
            except Exception as e:
                print(
                    f"   ⚠️  AST falló en {os.path.basename(file_path)} ({language}): {str(e)[:50]} — troceo por texto"
                )
                fallback += 1
        if doc_nodes is None:
            doc_nodes = _sentence_splitter().get_nodes_from_documents([doc])
            prose_count += 1
        nodes.extend(doc_nodes)

    print(f"   └─ {code_count} por AST, {prose_count} por texto ({fallback} con fallback) → {len(nodes)} chunks")
    emit_progress("chunking", chunks_generated=len(nodes), determinate=False)
    return nodes


def iter_node_batches(paths: list[str], context: SourceContext | None = None):
    for batch_number, batch in enumerate(iter_batches(paths), start=1):
        prepared = prepare_batch(batch, batch_number=batch_number, context=context)
        if prepared.nodes:
            yield prepared.nodes


def prepare_batch(
    paths: list[str],
    *,
    batch_number: int = 1,
    context: SourceContext | None = None,
) -> PreparedBatch:
    source_context = context or _default_source_context()
    loaded = load_docs_with_status(paths, source_context)
    prepared = PreparedBatch(failures=dict(loaded.failures))
    if not loaded.documents:
        return prepared
    print(f"   📦 Lote {batch_number}: {len(loaded.documents)} documentos, {len(paths)} archivos")
    documents_by_path: dict[str, list[Document]] = {}
    for document in loaded.documents:
        documents_by_path.setdefault(str(document.metadata.get("source_key") or ""), []).append(document)
    path_by_key = {source_context.source_key(path): path for path in loaded.loaded_paths}
    for source_key, documents in documents_by_path.items():
        path = path_by_key.get(source_key)
        if not path:
            continue
        try:
            nodes = build_nodes(documents)
        except Exception as exc:
            prepared.failures[path] = str(exc)[:300]
            print(f"   ⚠️  Error troceando {os.path.basename(path)}, se reintentará: {exc}")
            continue
        if not nodes:
            prepared.failures[path] = "chunking produced no nodes"
            continue
        prepared.nodes.extend(nodes)
        prepared.indexed_paths.add(path)
    return prepared


# ==================== MANIFIESTO (incremental) ====================
def _manifest_entry(value: object, context: SourceContext, relative: str) -> dict:
    if isinstance(value, dict):
        entry = dict(value)
    else:
        entry = {"legacy_fingerprint": value}
    entry.setdefault("schema_version", 1)
    entry["source_id"] = context.source_id
    entry["source_root"] = context.root
    entry["rel_path"] = relative
    return entry


def _migrate_manifest_for_context(raw: dict, context: SourceContext) -> dict:
    """Adopt pre-source manifests into the active root without losing peers.

    Legacy manifests used ``collection:relative-path`` and could only represent
    one root per collection.  The first run after this migration safely assigns
    those entries to the root being synchronized; modern entries already carrying
    ``source_id`` remain untouched.
    """
    migrated: dict[str, object] = {}
    active_prefix = f"{context.collection_id}:"
    for raw_key, value in raw.items():
        key = str(raw_key)
        if ":" not in key and context.collection_id == config.DEFAULT_COLLECTION_ID:
            key = f"{context.collection_id}:{key}"
        if not key.startswith(active_prefix):
            migrated[key] = value
            continue
        if isinstance(value, dict) and value.get("source_id"):
            migrated[key] = value
            continue
        relative = key[len(active_prefix) :]
        modern_key = context.source_key_for_relative(relative)
        migrated[modern_key] = _manifest_entry(value, context, relative)
    return migrated


def _expand_stored_manifest(raw: dict) -> dict:
    """Expand the backend-compatible on-disk envelope into source-flat state."""
    expanded: dict[str, object] = {}
    for raw_key, value in raw.items():
        key = str(raw_key)
        if not isinstance(value, dict) or value.get("manifest_schema") != MANIFEST_SCHEMA_VERSION:
            expanded[key] = value
            continue
        sources = value.get("sources")
        if not isinstance(sources, dict) or ":" not in key:
            expanded[key] = value
            continue
        collection_id, relative = key.split(":", 1)
        legacy = value.get("legacy")
        if legacy is not None:
            expanded[key] = legacy
        for source_id, entry in sources.items():
            safe_source_id = sanitize_collection_id(str(source_id), fallback="source")
            modern_key = f"{collection_id}:{safe_source_id}:{relative}"
            if isinstance(entry, dict):
                normalized = dict(entry)
                normalized["source_id"] = safe_source_id
                normalized.setdefault("rel_path", relative)
                expanded[modern_key] = normalized
            else:
                expanded[modern_key] = entry
    return expanded


def _manifest_for_storage(state: dict) -> dict:
    """Collapse source-flat state while keeping legacy backend delete keys.

    Backend source deletion historically removes ``collection:relative`` from
    the manifest.  Keeping that external key means old consumers still remove
    every colliding source they also remove from the index, while the nested
    source map lets the indexer synchronize roots independently.
    """
    stored: dict[str, object] = {}
    modern: dict[str, dict[str, object]] = {}
    for raw_key, value in state.items():
        key = str(raw_key)
        source_id = str(value.get("source_id") or "") if isinstance(value, dict) else ""
        if source_id and ":" in key:
            collection_id = key.split(":", 1)[0]
            modern_prefix = f"{collection_id}:{source_id}:"
            if key.startswith(modern_prefix):
                relative = key[len(modern_prefix) :]
                external_key = f"{collection_id}:{relative}"
                envelope = modern.setdefault(
                    external_key,
                    {"manifest_schema": MANIFEST_SCHEMA_VERSION, "sources": {}},
                )
                envelope["sources"][source_id] = value
                continue
        stored[key] = value
    for external_key, envelope in modern.items():
        if external_key in stored:
            envelope["legacy"] = stored.pop(external_key)
        stored[external_key] = envelope
    return stored


def read_manifest(context: SourceContext | None = None) -> dict:
    try:
        with open(config.MANIFEST_PATH, encoding="utf-8") as f:
            raw = json.load(f)
    except (OSError, ValueError):
        return {}
    if not isinstance(raw, dict):
        return {}
    canonical = {}
    for k, v in _expand_stored_manifest(raw).items():
        if ":" not in k:
            if COLLECTION_ID == config.DEFAULT_COLLECTION_ID:
                canonical[f"{COLLECTION_ID}:{k}"] = v
            else:
                canonical[k] = v
        else:
            canonical[k] = v
    return _migrate_manifest_for_context(canonical, context) if context is not None else canonical


def write_manifest(m: dict) -> None:
    os.makedirs(config.PERSIST_DIR, exist_ok=True)
    atomic_write_json(config.MANIFEST_PATH, _manifest_for_storage(m))


def _content_hash(path: str) -> str:
    digest = hashlib.blake2b(digest_size=32)
    with open(path, "rb") as stream:
        for block in iter(lambda: stream.read(_HASH_BLOCK_BYTES), b""):
            digest.update(block)
    return digest.hexdigest()


def current_state(paths: list[str], context: SourceContext | None = None) -> dict:
    """Return content-addressed, pipeline-versioned file fingerprints.

    ``mtime_ns`` and size keep diagnostics cheap, while the streaming content
    hash catches replacements that deliberately preserve both values.
    """
    state = {}
    source_context = context or _default_source_context()
    pipeline_version = _pipeline_version()
    for p in paths:
        try:
            stat = os.stat(p)
            relative = source_context.relative_path(p) if context is not None else _rel(p)
            key = source_context.source_key_for_relative(relative) if context is not None else _source_key(p)
            state[key] = {
                "schema_version": MANIFEST_SCHEMA_VERSION,
                "pipeline_version": pipeline_version,
                "hash_algorithm": FINGERPRINT_ALGORITHM,
                "content_hash": _content_hash(p),
                "mtime_ns": stat.st_mtime_ns,
                "size": stat.st_size,
                "source_id": source_context.source_id,
                "source_root": source_context.root,
                "rel_path": relative,
            }
        except (OSError, ValueError):
            pass
    return state


def _entry_belongs_to_source(key: str, value: object, context: SourceContext) -> bool:
    if isinstance(value, dict) and str(value.get("source_id") or "") == context.source_id:
        return key.startswith(f"{context.collection_id}:")
    return key.startswith(f"{context.collection_id}:{context.source_id}:")


def diff_manifest(
    old_state: dict,
    new_state: dict,
    rel_to_path: dict[str, str],
    context: SourceContext | None = None,
) -> tuple[list[str], list[str], list[str]]:
    new_files: list[str] = []
    changed: list[str] = []
    for key, path in rel_to_path.items():
        if key not in new_state:
            # File vanished/became unreadable between scanning and stat();
            # skip it here — it is handled by the `deleted` set below.
            continue
        if key in old_state:
            if old_state[key] != new_state[key]:
                changed.append(path)
        else:
            new_files.append(path)
    if APPEND_ONLY:
        deleted = []
    elif context is None:
        prefix = f"{COLLECTION_ID}:"
        deleted = [key for key in old_state if key.startswith(prefix) and key not in new_state]
    else:
        deleted = [
            key
            for key, value in old_state.items()
            if _entry_belongs_to_source(key, value, context) and key not in new_state
        ]
    return new_files, changed, deleted


def _relative_from_source_key(key: str, context: SourceContext | None = None) -> str:
    if context is not None:
        modern_prefix = f"{context.collection_id}:{context.source_id}:"
        if key.startswith(modern_prefix):
            return key[len(modern_prefix) :]
        collection_prefix = f"{context.collection_id}:"
        if key.startswith(collection_prefix):
            return key[len(collection_prefix) :]
    return key.split(":", 1)[1] if ":" in key else key


def remove_obsolete_nodes(
    index: VectorStoreIndex,
    changed: list[str],
    deleted: list[str],
    context: SourceContext | None = None,
) -> int:
    source_context = context
    source_keys_to_remove = (
        {source_context.source_key(path) for path in changed}
        if source_context is not None
        else {_source_key(path) for path in changed}
    ) | set(deleted)
    rels_to_remove = {_relative_from_source_key(key, source_context) for key in source_keys_to_remove}
    node_ids = []
    for nid, node in index.docstore.docs.items():
        metadata = node.metadata or {}
        source_key = str(metadata.get("source_key") or "")
        collection_id = str(metadata.get("collection_id") or config.DEFAULT_COLLECTION_ID)
        node_source_id = str(metadata.get("source_id") or "")
        if source_key in source_keys_to_remove:
            node_ids.append(nid)
            continue
        active_collection = source_context.collection_id if source_context is not None else COLLECTION_ID
        same_source = source_context is None or node_source_id in {"", source_context.source_id}
        legacy_key = f"{active_collection}:{metadata.get('rel_path', '')}"
        # Pre-source nodes can have either no source_key or the legacy
        # ``collection:relative`` key.  Only nodes without a source_id use this
        # fallback, so equal paths in another modern root remain untouched.
        legacy_match = not node_source_id and (not source_key or source_key == legacy_key)
        if (
            legacy_match
            and same_source
            and collection_id == active_collection
            and metadata.get("rel_path") in rels_to_remove
        ):
            node_ids.append(nid)
    if node_ids:
        index.delete_nodes(node_ids, delete_from_docstore=True)
    return len(node_ids)


@dataclass
class IndexUpdateResult:
    total_nodes: int = 0
    removed_nodes: int = 0
    indexed_paths: set[str] = field(default_factory=set)
    failures: dict[str, str] = field(default_factory=dict)


def apply_file_updates(
    index: VectorStoreIndex,
    paths: list[str],
    *,
    changed: set[str] | None = None,
    deleted: list[str] | None = None,
    context: SourceContext | None = None,
) -> IndexUpdateResult:
    """Extract, chunk and insert files without discarding good old chunks.

    Changed-file nodes are removed only after extraction/chunking succeeded.
    Failed files retain their previous nodes and remain absent/stale in the
    manifest so a later run retries them.
    """
    result = IndexUpdateResult()
    changed = changed or set()
    deleted = deleted or []
    if deleted:
        if context is None:
            result.removed_nodes += remove_obsolete_nodes(index, [], deleted)
        else:
            result.removed_nodes += remove_obsolete_nodes(index, [], deleted, context)
    if not paths:
        return result
    print("✂️  Troceando cambios...")
    for batch_number, batch in enumerate(iter_batches(paths), start=1):
        prepared = prepare_batch(batch, batch_number=batch_number, context=context)
        result.failures.update(prepared.failures)
        successful_changes = sorted(prepared.indexed_paths & changed)
        if successful_changes:
            if context is None:
                result.removed_nodes += remove_obsolete_nodes(index, successful_changes, [])
            else:
                result.removed_nodes += remove_obsolete_nodes(index, successful_changes, [], context)
        if not prepared.nodes:
            continue
        result.total_nodes += len(prepared.nodes)
        insert_node_batches(index, prepared.nodes)
        result.indexed_paths.update(prepared.indexed_paths)
    return result


def insert_files(index: VectorStoreIndex, paths: list[str], context: SourceContext | None = None) -> int:
    """Compatibility wrapper used by older integrations."""
    return apply_file_updates(index, paths, context=context).total_nodes


def _state_after_failures(
    old_state: dict,
    new_state: dict,
    failed_paths: set[str],
    context: SourceContext | None = None,
) -> dict:
    effective = dict(new_state)
    for path in failed_paths:
        key = context.source_key(path) if context is not None else _source_key(path)
        if key in old_state:
            effective[key] = old_state[key]
        else:
            effective.pop(key, None)
    return effective


def _merge_final_state(
    old_state: dict,
    new_state: dict,
    *,
    incremental: bool,
    context: SourceContext | None = None,
) -> dict:
    source_context = context
    if incremental and APPEND_ONLY:
        merged_state = dict(old_state)
        merged_state.update(new_state)
        return merged_state
    if source_context is None:
        prefix = f"{COLLECTION_ID}:"
        merged_state = {k: v for k, v in old_state.items() if not k.startswith(prefix)}
    else:
        merged_state = {
            k: v
            for k, v in old_state.items()
            if not _entry_belongs_to_source(k, v, source_context)
        }
    merged_state.update(new_state)
    return merged_state


def persist_final_state(
    old_state: dict,
    new_state: dict,
    *,
    incremental: bool,
    context: SourceContext | None = None,
) -> int:
    """Persist the manifest without clobbering entries from other collections.

    The manifest is global with ``{COLLECTION_ID}:{rel}`` keys, but ``new_state``
    only holds the active collection's files. Writing it verbatim would wipe every
    other collection's entries and force a full re-embed of them next time. So we
    keep foreign-collection keys untouched and only refresh the active prefix.
    """
    merged_state = _merge_final_state(
        old_state,
        new_state,
        incremental=incremental,
        context=context,
    )
    write_manifest(merged_state)
    return len(merged_state)


def run_incremental(
    old_state: dict,
    new_state: dict,
    rel_to_path: dict[str, str],
    context: SourceContext | None = None,
) -> int:
    from llama_index.core import StorageContext, load_index_from_storage

    source_context = context or _default_source_context()
    new_files, changed, deleted = diff_manifest(old_state, new_state, rel_to_path, source_context)
    if not (new_files or changed or deleted):
        print("\n✅ Todo al día — no hay cambios que indexar.")
        return 0

    print(f"\n🔄 Incremental: {len(new_files)} nuevos, {len(changed)} modificados, {len(deleted)} eliminados")
    print("📥 Cargando índice existente...")
    sc = StorageContext.from_defaults(persist_dir=config.PERSIST_DIR)
    index = load_index_from_storage(sc)

    update = apply_file_updates(
        index,
        new_files + changed,
        changed=set(changed),
        deleted=deleted,
        context=source_context,
    )
    if update.removed_nodes:
        print(f"   🗑️  {update.removed_nodes} chunks obsoletos eliminados")
    if update.failures:
        print(f"   ⚠️  {len(update.failures)} archivos conservaron su estado anterior y se reintentarán")
    effective_state = _state_after_failures(old_state, new_state, set(update.failures), source_context)
    merged_state = _merge_final_state(
        old_state,
        effective_state,
        incremental=True,
        context=source_context,
    )
    print("💾 Publicando generación atómica del índice...")
    publish_index_generation(
        index,
        _manifest_for_storage(merged_state),
        persist_dir=config.PERSIST_DIR,
        manifest_path=config.MANIFEST_PATH,
    )
    final_count = len(merged_state)
    print_summary(final_count, source_context)
    return 0


def _node_source_key(node, context: SourceContext | None = None) -> str | None:
    metadata = node.metadata or {}
    source_key = str(metadata.get("source_key") or "").strip()
    source_id = str(metadata.get("source_id") or "").strip()
    rel_path = str(metadata.get("rel_path") or "").strip()
    collection_id = str(metadata.get("collection_id") or config.DEFAULT_COLLECTION_ID)
    if context is not None and collection_id == context.collection_id and rel_path and not source_id:
        # Adopt nodes created before source roots had identities.
        return context.source_key_for_relative(rel_path)
    if source_key:
        return source_key
    if not rel_path:
        return None
    if source_id:
        return f"{collection_id}:{source_id}:{rel_path}"
    return f"{collection_id}:{rel_path}"


def run_manifest_recovery(
    new_state: dict,
    rel_to_path: dict[str, str],
    context: SourceContext | None = None,
) -> int:
    """Recover a missing/corrupt manifest without replacing other collections."""
    from llama_index.core import StorageContext, load_index_from_storage

    print("\n🛟 Índice existente sin manifiesto válido — recuperación segura")
    storage_context = StorageContext.from_defaults(persist_dir=config.PERSIST_DIR)
    existing = load_index_from_storage(storage_context)
    source_context = context or _default_source_context()
    node_context = context
    existing_keys = {
        key
        for node in existing.docstore.docs.values()
        if (key := _node_source_key(node, node_context)) is not None
    }
    active_prefix = (
        f"{source_context.collection_id}:{source_context.source_id}:"
        if context is not None
        else f"{COLLECTION_ID}:"
    )
    deleted = sorted(key for key in existing_keys if key.startswith(active_prefix) and key not in new_state)
    paths = list(rel_to_path.values())
    update = apply_file_updates(
        existing,
        paths,
        changed=set(paths),
        deleted=deleted,
        context=source_context,
    )

    recovered: dict[str, dict] = {}
    successful_keys = {
        key for key, path in rel_to_path.items() if path in update.indexed_paths
    }
    for node in existing.docstore.docs.values():
        key = _node_source_key(node, node_context)
        if not key:
            continue
        if key in successful_keys and key in new_state:
            recovered[key] = new_state[key]
        else:
            # A deliberately non-matching fingerprint forces verification when
            # that collection is indexed next, while preserving its nodes now.
            metadata = node.metadata or {}
            recovered_entry = {"unverified": True}
            if metadata.get("source_id"):
                recovered_entry.update(
                    {
                        "schema_version": MANIFEST_SCHEMA_VERSION,
                        "pipeline_version": _pipeline_version(),
                        "source_id": metadata.get("source_id"),
                        "source_root": metadata.get("source_root"),
                        "rel_path": metadata.get("rel_path"),
                    }
                )
            recovered[key] = recovered_entry
    print("💾 Publicando generación recuperada del índice...")
    publish_index_generation(
        existing,
        _manifest_for_storage(recovered),
        persist_dir=config.PERSIST_DIR,
        manifest_path=config.MANIFEST_PATH,
    )
    if update.failures:
        print(f"   ⚠️  {len(update.failures)} archivos se conservaron y se reintentarán")
    print_summary(len(recovered), source_context)
    return 0


def run_full_index(
    paths: list[str],
    new_state: dict,
    context: SourceContext | None = None,
) -> int:
    source_context = context or _default_source_context()
    print("\n🆕 Indexado completo (primera vez)")
    if not paths:
        print("❌ No se encontraron documentos para indexar.")
        return 1
    print("✂️  Troceando (chunking consciente del lenguaje)...")
    index = None
    total_nodes = 0
    indexed_paths: set[str] = set()
    failures: dict[str, str] = {}
    for batch_number, batch in enumerate(iter_batches(paths), start=1):
        prepared = prepare_batch(batch, batch_number=batch_number, context=source_context)
        failures.update(prepared.failures)
        nodes = prepared.nodes
        if not nodes:
            continue
        indexed_paths.update(prepared.indexed_paths)
        total_nodes += len(nodes)
        index = insert_node_batches(index, nodes, initialize=index is None)
    if index is None:
        print("❌ No se pudieron generar chunks para indexar.")
        return 1
    successful_state = {
        source_context.source_key(path): new_state[source_context.source_key(path)]
        for path in indexed_paths
        if source_context.source_key(path) in new_state
    }
    print("💾 Publicando primera generación atómica del índice...")
    publish_index_generation(
        index,
        _manifest_for_storage(successful_state),
        persist_dir=config.PERSIST_DIR,
        manifest_path=config.MANIFEST_PATH,
    )
    final_count = len(successful_state)
    if failures:
        print(f"   ⚠️  {len(failures)} archivos no se marcaron y se reintentarán")
    print_summary(final_count, source_context)
    return 0


def print_summary(final_count: int, context: SourceContext | None = None) -> None:
    source_context = context or _default_source_context()
    print("\n✅ Indexado completado")
    print(f"📚 Colección: {source_context.collection_name} ({source_context.collection_id})")
    print(f"🗂️  Fuente: {source_context.project_name} ({source_context.source_id})")
    print(f"📦 {config.PERSIST_DIR}  ·  {final_count} archivos en el índice")
    print("═" * 45)


def run_index(root: str | None = None) -> int:
    root = root or config.PROJECTS_DIRS[0]
    source_context = SourceContext.create(root)
    print("\n🧠 TrinaxAI — Indexador de Documentos")
    print("═" * 45)
    if not os.path.isdir(root):
        print(f"❌ Directorio no encontrado: {root}")
        return 1
    lock_timeout = _env_int("TRINAXAI_INDEX_LOCK_TIMEOUT", 3600, minimum=1, maximum=86400)
    lock_path = os.path.join(config.PERSIST_DIR, ".indexing.lock")
    print("🔒 Esperando turno exclusivo del índice...", flush=True)
    try:
        with exclusive_process_lock(lock_path, timeout=lock_timeout):
            recovery = recover_interrupted_transaction(config.PERSIST_DIR, config.MANIFEST_PATH)
            if recovery == "rolled_back":
                print("🛟 Se restauró la generación anterior tras una indexación interrumpida.")
            elif recovery == "committed":
                print("🧹 Se confirmó una generación ya publicada y se limpió su transacción.")
            ensure_embed_settings()
            print(f"📂 Recorriendo: {source_context.root}")
            paths = collect_files(source_context.root)
            rel_to_path = {source_context.source_key(p): p for p in paths}
            old_state = read_manifest(source_context)
            new_state = current_state(paths, source_context)
            print(f"   └─ {len(paths)} archivos candidatos")

            index_exists = os.path.exists(os.path.join(config.PERSIST_DIR, "docstore.json"))
            if index_exists and old_state:
                return run_incremental(old_state, new_state, rel_to_path, source_context)
            if index_exists:
                return run_manifest_recovery(new_state, rel_to_path, source_context)
            return run_full_index(paths, new_state, source_context)
    except TimeoutError as exc:
        print(f"❌ {exc}")
        return 2
    except (OSError, RuntimeError, ValueError) as exc:
        print(f"❌ No se pudo publicar/recuperar el índice: {exc}")
        return 3


# ==================== MAIN ====================
if __name__ == "__main__":
    sys.exit(run_index())
