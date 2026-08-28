from __future__ import annotations

import sys
from io import BytesIO
from threading import BoundedSemaphore
from types import SimpleNamespace
from zipfile import ZIP_DEFLATED, ZipFile

import pytest
from fastapi import HTTPException
from starlette.testclient import TestClient

import rag_api
from app.services import document_service


def test_extracts_plain_text_through_api() -> None:
    response = TestClient(rag_api.app, client=("127.0.0.1", 50000)).post(
        "/documents/extract",
        files={"file": ("notes.txt", "Información local".encode(), "text/plain")},
    )

    assert response.status_code == 200
    assert response.json() == {
        "ok": True,
        "name": "notes.txt",
        "text": "Información local",
        "chars": 17,
        "truncated": False,
    }


def test_document_slot_is_released_when_extraction_fails(monkeypatch) -> None:
    slots = BoundedSemaphore(1)
    monkeypatch.setattr(document_service, "_document_slots", slots)
    monkeypatch.setattr(
        document_service,
        "_extract_document_text",
        lambda *_args: (_ for _ in ()).throw(ValueError("broken document")),
    )

    response = TestClient(rag_api.app, client=("127.0.0.1", 50000), raise_server_exceptions=False).post(
        "/documents/extract",
        files={"file": ("broken.txt", b"content", "text/plain")},
    )

    assert response.status_code == 500
    assert slots.acquire(blocking=False)
    slots.release()


def test_extracts_pptx_text() -> None:
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "TrinaxAI slides"
    slide.placeholders[1].text = "Readable presentation content"
    data = BytesIO()
    presentation.save(data)

    text = rag_api._extract_document_text("deck.pptx", data.getvalue())

    assert "TrinaxAI slides" in text
    assert "Readable presentation content" in text


def test_extracts_docx_text_and_removes_temporary_file() -> None:
    data = BytesIO()
    with ZipFile(data, "w", ZIP_DEFLATED) as archive:
        archive.writestr(
            "word/document.xml",
            (
                '<?xml version="1.0" encoding="UTF-8"?>'
                '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
                "<w:body><w:p><w:r><w:t>Release plan</w:t></w:r></w:p>"
                "<w:p><w:r><w:t>Production checklist</w:t></w:r></w:p></w:body></w:document>"
            ),
        )

    text = rag_api._extract_document_text("plan.docx", data.getvalue())

    assert "Release plan" in text
    assert "Production checklist" in text


def test_rejects_archives_that_expand_past_the_parser_budget(monkeypatch) -> None:
    data = BytesIO()
    with ZipFile(data, "w", ZIP_DEFLATED) as archive:
        archive.writestr("word/document.xml", "1234")
    monkeypatch.setattr(document_service, "_ARCHIVE_MAX_UNCOMPRESSED_BYTES", 3)

    with pytest.raises(HTTPException) as raised:
        document_service._extract_document_text("plan.docx", data.getvalue())

    assert raised.value.status_code == 413


def test_rejects_invalid_office_archive() -> None:
    with pytest.raises(HTTPException) as raised:
        document_service._extract_document_text("broken.docx", b"not a zip")

    assert raised.value.status_code == 422


def test_rejects_office_archive_with_too_many_members(monkeypatch) -> None:
    data = BytesIO()
    with ZipFile(data, "w", ZIP_DEFLATED) as archive:
        archive.writestr("one", "1")
        archive.writestr("two", "2")
    monkeypatch.setattr(document_service, "_ARCHIVE_MAX_MEMBERS", 1)

    with pytest.raises(HTTPException) as raised:
        document_service._extract_document_text("plan.docx", data.getvalue())

    assert raised.value.status_code == 413


def test_extracts_pdf_pages_and_optional_ocr_fallback(monkeypatch) -> None:
    pages = [
        SimpleNamespace(extract_text=lambda: "First page"),
        SimpleNamespace(extract_text=lambda: ""),
    ]
    monkeypatch.setitem(sys.modules, "pypdf", SimpleNamespace(PdfReader=lambda _stream: SimpleNamespace(pages=pages)))
    monkeypatch.setattr(document_service.config, "TRINAXAI_OCR", False)

    assert document_service._extract_pdf_text(b"pdf") == "[Page 1]\nFirst page"


def test_pdf_ocr_replaces_empty_page_text_when_enabled(monkeypatch) -> None:
    page = SimpleNamespace(extract_text=lambda: "")
    monkeypatch.setitem(
        sys.modules,
        "pypdf",
        SimpleNamespace(PdfReader=lambda _stream: SimpleNamespace(pages=[page])),
    )
    monkeypatch.setitem(
        sys.modules,
        "pdf2image",
        SimpleNamespace(convert_from_bytes=lambda *_args, **_kwargs: [object()]),
    )
    monkeypatch.setitem(
        sys.modules,
        "pytesseract",
        SimpleNamespace(image_to_string=lambda *_args, **_kwargs: "Scanned production release checklist"),
    )
    monkeypatch.setattr(document_service.config, "TRINAXAI_OCR", True)

    assert "Scanned production release checklist" in document_service._extract_pdf_text(b"pdf")


def test_extracts_xlsx_text() -> None:
    from openpyxl import Workbook

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Budget"
    sheet.append(["Item", "Amount"])
    sheet.append(["Hosting", 25])
    data = BytesIO()
    workbook.save(data)

    text = rag_api._extract_document_text("budget.xlsx", data.getvalue())

    assert "[Sheet: Budget]" in text
    assert "Hosting\t25" in text


def test_extracts_rtf_text() -> None:
    text = rag_api._extract_document_text("notes.rtf", rb"{\rtf1\ansi TrinaxAI notes}")
    assert "TrinaxAI notes" in text


def test_extracts_odt_text() -> None:
    from odf.opendocument import OpenDocumentText
    from odf.text import P

    document = OpenDocumentText()
    document.text.addElement(P(text="OpenDocument content for TrinaxAI"))
    data = BytesIO()
    document.save(data)

    text = rag_api._extract_document_text("notes.odt", data.getvalue())

    assert "OpenDocument content for TrinaxAI" in text


def test_legacy_office_conversion_requires_libreoffice(monkeypatch) -> None:
    monkeypatch.setattr(document_service.shutil, "which", lambda _name: None)

    with pytest.raises(HTTPException) as raised:
        document_service._convert_office_bytes(b"legacy", ".doc", ".docx")

    assert raised.value.status_code == 501


def test_legacy_office_conversion_returns_generated_container(monkeypatch) -> None:
    monkeypatch.setattr(document_service.shutil, "which", lambda _name: "/usr/bin/libreoffice")

    def convert(command, **_kwargs):
        output_dir = command[command.index("--outdir") + 1]
        with open(f"{output_dir}/document.docx", "wb") as stream:
            stream.write(b"converted")
        return SimpleNamespace(returncode=0, stdout="", stderr="")

    monkeypatch.setattr(document_service.subprocess, "run", convert)

    assert document_service._convert_office_bytes(b"legacy", ".doc", ".docx") == b"converted"


@pytest.mark.parametrize(
    ("extractor", "payload"),
    [
        (document_service._extract_pdf_text, b"not a pdf"),
        (document_service._extract_docx_text, b"not a docx"),
        (document_service._extract_pptx_text, b"not a pptx"),
        (document_service._extract_xlsx_text, b"not an xlsx"),
        (document_service._extract_odf_text, b"not an odf"),
    ],
)
def test_corrupt_structured_documents_return_unprocessable(extractor, payload) -> None:
    with pytest.raises(HTTPException) as raised:
        extractor(payload)
    assert raised.value.status_code == 422


@pytest.mark.parametrize(
    ("filename", "target_extension", "extractor_name"),
    [
        ("legacy.doc", ".docx", "_extract_docx_text"),
        ("slides.ppt", ".pptx", "_extract_pptx_text"),
        ("sheet.xls", ".xlsx", "_extract_xlsx_text"),
    ],
)
def test_legacy_formats_dispatch_through_conversion(
    monkeypatch,
    filename,
    target_extension,
    extractor_name,
) -> None:
    monkeypatch.setattr(
        document_service,
        "_convert_office_bytes",
        lambda data, source, target: b"converted" if data == b"legacy" and target == target_extension else b"",
    )
    monkeypatch.setattr(document_service, extractor_name, lambda data: f"read {data.decode()}")

    assert document_service._extract_document_text(filename, b"legacy") == "read converted"


def test_document_api_rejects_empty_and_oversized_files(monkeypatch) -> None:
    client = TestClient(rag_api.app, client=("127.0.0.1", 50000))
    empty = client.post("/documents/extract", files={"file": ("empty.txt", b"", "text/plain")})
    assert empty.status_code == 400

    monkeypatch.setattr(document_service, "DOC_EXTRACT_MAX_BYTES", 3)
    oversized = client.post("/documents/extract", files={"file": ("large.txt", b"four", "text/plain")})
    assert oversized.status_code == 413


def test_document_api_truncates_large_extracted_text(monkeypatch) -> None:
    monkeypatch.setattr(document_service, "DOC_EXTRACT_MAX_CHARS", 4)

    response = TestClient(rag_api.app, client=("127.0.0.1", 50000)).post(
        "/documents/extract",
        files={"file": ("notes.txt", b"abcdef", "text/plain")},
    )

    assert response.status_code == 200
    assert response.json()["text"] == "abcd"
    assert response.json()["truncated"] is True


def test_memory_fallback_never_persists_model_error(monkeypatch, tmp_path) -> None:
    monkeypatch.setattr(rag_api, "_memory_load", lambda: {"memories": [{"text": "Prefiere respuestas breves"}]})
    monkeypatch.setattr(rag_api.config, "PERSIST_DIR", str(tmp_path))
    monkeypatch.setattr(rag_api, "get_llm", lambda *_args, **_kwargs: (_ for _ in ()).throw(RuntimeError("offline")))

    result = rag_api._memory_refresh_sync(rag_api.MemoryRefreshRequest())

    assert result["summary"] == "Prefiere respuestas breves"
    assert "LLM unavailable" not in result["summary"]
