"""File discovery and document extraction for the incremental indexer."""

from __future__ import annotations

import json
import os
import shutil
import subprocess
import tempfile
import zipfile
from dataclasses import dataclass
from email import policy
from email.parser import BytesParser
from html.parser import HTMLParser
from typing import TYPE_CHECKING

if TYPE_CHECKING:
    from llama_index.core.schema import Document

import config
from trinaxai_core import sanitize_collection_id, source_id_for_root

EXTRACTOR_EXTS = {".pdf", ".docx"}
TEXT_ENCODINGS = ("utf-8-sig", "utf-8", "cp1252", "latin-1")
_TEXT_SAMPLE_BYTES = 8192
_EPUB_TEXT_LIMIT = 20 * 1024 * 1024
_SENSITIVE_NAMES = {".env", ".netrc", ".npmrc", ".pypirc", "credentials.json", "secrets.json", "id_rsa", "id_ed25519"}
_SENSITIVE_EXTENSIONS = {".key", ".pem", ".p12", ".pfx"}
COLLECTION_ID = sanitize_collection_id(
    os.getenv("TRINAXAI_COLLECTION_ID", config.DEFAULT_COLLECTION_ID),
    fallback=config.DEFAULT_COLLECTION_ID,
)
COLLECTION_NAME = (
    os.getenv("TRINAXAI_COLLECTION_NAME", config.DEFAULT_COLLECTION_NAME).strip() or config.DEFAULT_COLLECTION_NAME
)


class _HTMLTextExtractor(HTMLParser):
    """Small dependency-free HTML-to-text extractor."""

    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.parts: list[str] = []
        self._ignored_depth = 0

    def handle_starttag(self, tag: str, attrs) -> None:
        if tag in {"script", "style", "noscript", "svg"}:
            self._ignored_depth += 1
        elif tag in {"p", "div", "br", "li", "h1", "h2", "h3", "h4", "tr"}:
            self.parts.append("\n")

    def handle_endtag(self, tag: str) -> None:
        if tag in {"script", "style", "noscript", "svg"} and self._ignored_depth:
            self._ignored_depth -= 1
        elif tag in {"p", "div", "li", "h1", "h2", "h3", "h4", "tr"}:
            self.parts.append("\n")

    def handle_data(self, data: str) -> None:
        if not self._ignored_depth and data.strip():
            self.parts.append(data.strip())

    def text(self) -> str:
        lines = (" ".join(line.split()) for line in " ".join(self.parts).splitlines())
        return "\n".join(line for line in lines if line).strip()


def _html_to_text(value: str) -> str:
    parser = _HTMLTextExtractor()
    parser.feed(value)
    parser.close()
    return parser.text()


@dataclass(frozen=True)
class SourceContext:
    """Identity and path policy for one independently synchronized source root."""

    root: str
    source_id: str
    collection_id: str
    collection_name: str
    project_name: str

    @classmethod
    def create(
        cls,
        root: str,
        *,
        source_id: str | None = None,
        collection_id: str | None = None,
        collection_name: str | None = None,
    ) -> "SourceContext":
        canonical_root = os.path.realpath(os.path.abspath(os.path.expanduser(root)))
        basename = os.path.basename(canonical_root.rstrip(os.sep)) or "root"
        explicit_id = source_id or os.getenv("TRINAXAI_SOURCE_ID")
        safe_source_id = source_id_for_root(canonical_root, explicit_id=explicit_id)
        return cls(
            root=canonical_root,
            source_id=safe_source_id,
            collection_id=sanitize_collection_id(
                collection_id or COLLECTION_ID,
                fallback=config.DEFAULT_COLLECTION_ID,
            ),
            collection_name=(collection_name or COLLECTION_NAME).strip() or config.DEFAULT_COLLECTION_NAME,
            project_name=basename,
        )

    def relative_path(self, path: str) -> str:
        absolute = os.path.realpath(os.path.abspath(path))
        try:
            if os.path.commonpath([self.root, absolute]) != self.root:
                raise ValueError(f"Path is outside source root: {path}")
        except ValueError as exc:
            raise ValueError(f"Path is outside source root: {path}") from exc
        return os.path.relpath(absolute, self.root).replace("\\", "/")

    def source_key_for_relative(self, relative: str) -> str:
        clean_relative = relative.replace("\\", "/").lstrip("/")
        return f"{self.collection_id}:{self.source_id}:{clean_relative}"

    def source_key(self, path: str) -> str:
        return self.source_key_for_relative(self.relative_path(path))


def default_source_context() -> SourceContext:
    return SourceContext.create(config.PROJECTS_DIRS[0])


def collect_files(root: str) -> list[str]:
    """Walk ``root`` while pruning dependencies, hidden folders and secrets."""
    allowed = {extension.lower() for extension in config.REQUIRED_EXTS}
    allowed_names = {
        "dockerfile",
        "makefile",
        "readme",
        "license",
        "changelog",
        "contributing",
        "gemfile",
        "procfile",
    }
    files: list[str] = []
    skipped_big = 0
    file_count = 0
    for dirpath, dirnames, filenames in os.walk(root):
        dirnames[:] = [
            directory
            for directory in dirnames
            if directory not in config.EXCLUDE_DIR_NAMES
            and not directory.startswith(".")
            and not directory.endswith((".egg-info", ".dist-info"))
        ]
        for filename in filenames:
            lower_name = filename.lower()
            if (
                lower_name in _SENSITIVE_NAMES
                or lower_name.startswith(".env.")
                or os.path.splitext(lower_name)[1] in _SENSITIVE_EXTENSIONS
            ):
                continue
            if filename.startswith(".") and filename.lower() not in allowed_names:
                continue
            full_path = os.path.join(dirpath, filename)
            if os.path.islink(full_path):
                continue
            try:
                if os.path.getsize(full_path) > config.max_file_bytes(full_path):
                    skipped_big += 1
                    continue
            except OSError:
                continue
            known_type = filename.lower() in allowed_names or os.path.splitext(filename)[1].lower() in allowed
            if not known_type and not is_probably_text_file(full_path):
                continue
            files.append(full_path)
            file_count += 1
            if file_count % 5000 == 0:
                print(f"   📂 {file_count} archivos encontrados...", flush=True)
    if file_count >= 5000:
        print(f"   📂 {file_count} archivos encontrados en total")
    if skipped_big:
        print(f"   ⏭️  {skipped_big} archivos omitidos por tamaño (sobre el límite configurado para su tipo de archivo)")
    return files


def is_probably_text_file(path: str) -> bool:
    try:
        with open(path, "rb") as stream:
            sample = stream.read(_TEXT_SAMPLE_BYTES)
    except OSError:
        return False
    if not sample:
        return True
    if b"\x00" in sample:
        return False
    try:
        decoded = sample.decode("utf-8")
    except UnicodeDecodeError:
        decoded = sample.decode("cp1252", errors="replace")
    controls = sum(1 for char in decoded if ord(char) < 32 and char not in {"\n", "\r", "\t", "\f", "\b"})
    replacements = decoded.count("\ufffd")
    return (controls + replacements) / max(1, len(decoded)) < 0.02


def relative_path(path: str, context: SourceContext | None = None) -> str:
    if context is not None:
        return context.relative_path(path)
    try:
        return os.path.relpath(path, config.PROJECTS_DIRS[0]).replace("\\", "/")
    except ValueError:
        return path.replace("\\", "/")


def source_key(path: str, context: SourceContext | None = None) -> str:
    if context is not None:
        return context.source_key(path)
    return f"{COLLECTION_ID}:{relative_path(path)}"


def decode_text_bytes(data: bytes) -> str:
    if data.startswith((b"\xff\xfe", b"\xfe\xff")) or b"\x00" in data[:200]:
        try:
            return data.decode("utf-16")
        except UnicodeDecodeError:
            pass
    for encoding in TEXT_ENCODINGS:
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def _load_text_document(path: str) -> Document:
    from llama_index.core.schema import Document

    with open(path, "rb") as stream:
        return Document(text=decode_text_bytes(stream.read()), metadata={"file_path": path})


def document(path: str, text: str) -> Document:
    from llama_index.core.schema import Document

    cleaned = text.strip()
    if not cleaned:
        raise ValueError("the file contains no extractable text")
    return Document(text=cleaned, metadata={"file_path": path})


def _load_html_document(path: str) -> Document:
    with open(path, "rb") as stream:
        return document(path, _html_to_text(decode_text_bytes(stream.read())))


def _load_notebook_document(path: str) -> Document:
    with open(path, encoding="utf-8") as stream:
        notebook = json.load(stream)
    parts: list[str] = []
    for number, cell in enumerate(notebook.get("cells") or [], start=1):
        if not isinstance(cell, dict):
            continue
        cell_type = str(cell.get("cell_type") or "cell")
        source = cell.get("source") or ""
        text = "".join(source) if isinstance(source, list) else str(source)
        if text.strip():
            parts.append(f"[{cell_type.title()} cell {number}]\n{text.strip()}")
    return document(path, "\n\n".join(parts))


def _load_email_document(path: str) -> Document:
    with open(path, "rb") as stream:
        message = BytesParser(policy=policy.default).parse(stream)
    headers = [f"{name}: {message.get(name)}" for name in ("Subject", "From", "To", "Date") if message.get(name)]
    bodies: list[str] = []
    parts = message.walk() if message.is_multipart() else [message]
    for part in parts:
        if part.get_content_disposition() == "attachment":
            continue
        content_type = part.get_content_type()
        if content_type not in {"text/plain", "text/html"}:
            continue
        try:
            value = part.get_content()
        except Exception:
            value = decode_text_bytes(part.get_payload(decode=True) or b"")
        bodies.append(_html_to_text(value) if content_type == "text/html" else value.strip())
    return document(path, "\n".join(headers) + "\n\n" + "\n\n".join(bodies))


def _load_epub_document(path: str) -> Document:
    sections: list[str] = []
    extracted_bytes = 0
    with zipfile.ZipFile(path) as archive:
        for info in archive.infolist():
            extension = os.path.splitext(info.filename.lower())[1]
            if extension not in {".html", ".htm", ".xhtml", ".xml", ".ncx"} or info.file_size > _EPUB_TEXT_LIMIT:
                continue
            extracted_bytes += info.file_size
            if extracted_bytes > _EPUB_TEXT_LIMIT:
                raise ValueError("EPUB expanded text exceeds the safe extraction limit")
            value = _html_to_text(decode_text_bytes(archive.read(info)))
            if value:
                sections.append(f"[{info.filename}]\n{value}")
    return document(path, "\n\n".join(sections))


def _load_extracted_documents(path: str) -> list[Document]:
    from llama_index.core import SimpleDirectoryReader

    return SimpleDirectoryReader(
        input_files=[path],
        exclude_hidden=False,
        exclude=config.EXCLUDE_PATTERNS,
        encoding="utf-8",
        errors="replace",
        raise_on_error=False,
    ).load_data()


def emit_progress(phase: str, **values: object) -> None:
    print("TRINAXAI_PROGRESS " + json.dumps({"phase": phase, **values}, ensure_ascii=False), flush=True)


def _load_pdf_documents(path: str) -> list[Document]:
    from llama_index.core.schema import Document
    from pypdf import PdfReader

    reader = PdfReader(path, strict=False)
    total = len(reader.pages)
    emit_progress("extracting", pages_total=total, pages_processed=0, determinate=bool(total))
    documents: list[Document] = []
    for page_number, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        if text:
            documents.append(
                Document(text=f"[Page {page_number}]\n{text}", metadata={"file_path": path, "page": page_number})
            )
        emit_progress("extracting", pages_total=total, pages_processed=page_number, determinate=True)
    if not documents:
        raise ValueError("PDF contains no extractable text; OCR may be required")
    return documents


def _load_pptx_document(path: str) -> Document:
    from llama_index.core.schema import Document
    from pptx import Presentation

    presentation = Presentation(path)
    slides: list[str] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            text = ""
            if getattr(shape, "has_text_frame", False):
                text = shape.text or ""
            elif getattr(shape, "has_table", False):
                rows = []
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        rows.append(" | ".join(cells))
                text = "\n".join(rows)
            if text.strip():
                parts.append(text.strip())
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text or ""
                if notes.strip():
                    parts.append(f"Notes:\n{notes.strip()}")
        except Exception:
            pass
        if parts:
            slides.append(f"[Slide {slide_number}]\n" + "\n\n".join(parts))
    return Document(text="\n\n".join(slides).strip(), metadata={"file_path": path})


def _load_xlsx_document(path: str) -> Document:
    from llama_index.core.schema import Document
    from openpyxl import load_workbook

    workbook = load_workbook(path, read_only=True, data_only=True)
    sheets: list[str] = []
    try:
        for sheet in workbook.worksheets:
            rows: list[str] = []
            for row in sheet.iter_rows(values_only=True):
                values = [str(value).strip() if value is not None else "" for value in row]
                while values and not values[-1]:
                    values.pop()
                if any(values):
                    rows.append("\t".join(values))
            if rows:
                sheets.append(f"[Sheet: {sheet.title}]\n" + "\n".join(rows))
    finally:
        workbook.close()
    return Document(text="\n\n".join(sheets).strip(), metadata={"file_path": path})


def _load_rtf_document(path: str) -> Document:
    from llama_index.core.schema import Document
    from striprtf.striprtf import rtf_to_text

    with open(path, "rb") as stream:
        text = rtf_to_text(decode_text_bytes(stream.read()))
    return Document(text=text.strip(), metadata={"file_path": path})


def _load_odf_document(path: str) -> Document:
    from llama_index.core.schema import Document
    from odf import teletype
    from odf.opendocument import load
    from odf.table import Table, TableCell, TableRow
    from odf.text import H, P

    document_value = load(path)
    parts: list[str] = []
    for node_type in (H, P):
        for node in document_value.getElementsByType(node_type):
            text = teletype.extractText(node).strip()
            if text:
                parts.append(text)
    for table in document_value.getElementsByType(Table):
        rows: list[str] = []
        for row in table.getElementsByType(TableRow):
            cells = [teletype.extractText(cell).strip() for cell in row.getElementsByType(TableCell)]
            while cells and not cells[-1]:
                cells.pop()
            if any(cells):
                rows.append("\t".join(cells))
        if rows:
            parts.append("\n".join(rows))
    return Document(text="\n".join(parts).strip(), metadata={"file_path": path})


def _load_converted_office_document(path: str, target_ext: str, loader=None) -> list[Document]:
    executable = shutil.which("libreoffice") or shutil.which("soffice")
    if not executable:
        raise RuntimeError(f"{os.path.splitext(path)[1].upper()} requires LibreOffice")
    with tempfile.TemporaryDirectory(prefix="trinaxai-index-office-") as directory:
        result = subprocess.run(
            [executable, "--headless", "--convert-to", target_ext.lstrip("."), "--outdir", directory, path],
            capture_output=True,
            text=True,
            timeout=90,
        )
        target = os.path.join(directory, f"{os.path.splitext(os.path.basename(path))[0]}{target_ext}")
        if result.returncode != 0 or not os.path.isfile(target):
            detail = (result.stderr or result.stdout or "conversion produced no output").strip()
            raise RuntimeError(f"Office conversion failed: {detail[:180]}")
        documents = (loader or load_file_documents)(target)
        for loaded in documents:
            loaded.metadata["file_path"] = path
        return documents


def load_file_documents(path: str) -> list[Document]:
    extension = os.path.splitext(path)[1].lower()
    if extension == ".pdf":
        return _load_pdf_documents(path)
    if extension in {".html", ".htm", ".xhtml"}:
        return [_load_html_document(path)]
    if extension == ".ipynb":
        return [_load_notebook_document(path)]
    if extension == ".eml":
        return [_load_email_document(path)]
    if extension == ".epub":
        return [_load_epub_document(path)]
    if extension == ".pptx":
        return [_load_pptx_document(path)]
    if extension == ".xlsx":
        return [_load_xlsx_document(path)]
    if extension == ".rtf":
        return [_load_rtf_document(path)]
    if extension in {".odt", ".ods", ".odp"}:
        return [_load_odf_document(path)]
    if extension == ".doc":
        return _load_converted_office_document(path, ".docx")
    if extension == ".ppt":
        return _load_converted_office_document(path, ".pptx")
    if extension == ".xls":
        return _load_converted_office_document(path, ".xlsx")
    if extension in EXTRACTOR_EXTS:
        return _load_extracted_documents(path)
    return [_load_text_document(path)]


def load_file_documents_result(path: str) -> tuple[str, list[Document], Exception | None]:
    try:
        return path, load_file_documents(path), None
    except Exception as exc:
        return path, [], exc
